from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT="Guia_despliegue_Symfony_React_Hostinger.pptx"
N=RGBColor(12,19,35); N2=RGBColor(23,34,56); W=RGBColor(255,255,255)
M=RGBColor(177,187,204); C=RGBColor(54,211,221); O=RGBColor(255,158,88)
G=RGBColor(83,214,142); R=RGBColor(255,105,120); I=RGBColor(27,32,43)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def box(s,x,y,w,h,color,round=False,line=None):
 q=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); q.fill.solid(); q.fill.fore_color.rgb=color; q.line.color.rgb=line or color; return q
def txt(s,t,x,y,w,h,z=18,c=W,b=False,font="Aptos",align=PP_ALIGN.LEFT):
 q=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); f=q.text_frame; f.clear(); f.margin_left=f.margin_right=Inches(.04); p=f.paragraphs[0]; p.text=t; p.alignment=align
 for r in p.runs: r.font.name=font; r.font.size=Pt(z); r.font.bold=b; r.font.color.rgb=c
 return q
def base(sec,num,title,sub=""):
 s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=N
 ac=O if sec=="FRONTEND" else C if sec=="BACKEND" else G; box(s,0,0,.13,7.5,ac); txt(s,sec,.55,.35,3,.25,10,ac,True); txt(s,f"{num:02d}",12.15,.35,.55,.25,10,M,True,align=PP_ALIGN.RIGHT); txt(s,title,.55,.83,11.7,.55,27,W,True)
 if sub: txt(s,sub,.58,1.48,11.6,.4,13,M)
 txt(s,"CINE APP · HOSTINGER COMPARTIDO",.58,7.12,4,.18,8,M,True); return s
def card(s,title,body,x,y,w,h,ac=C,tag=""):
 box(s,x,y,w,h,N2,True,RGBColor(42,57,80)); box(s,x,y,.07,h,ac,True); txt(s,title,x+.25,y+.2,w-.45,.32,15,W,True); txt(s,body,x+.25,y+.7,w-.45,h-.85,12.5,M)
 if tag: box(s,x+w-1.15,y+.18,.85,.3,ac,True); txt(s,tag,x+w-1.1,y+.23,.75,.14,8,N,True,align=PP_ALIGN.CENTER)

# portada
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=N; box(s,0,0,.14,7.5,C)
txt(s,"CINE APP / GUÍA DE PRODUCCIÓN",.7,.65,5,.3,11,C,True); txt(s,"De localhost a\nHostinger compartido",.68,1.45,7.2,1.55,35,W,True); txt(s,"Symfony 7.4 + React/Vite + MySQL",.72,3.4,6.4,.4,18,M)
card(s,"BACKEND","PHP 8.2 · Symfony · Doctrine\nMySQL · variables seguras",8.6,1.25,3.7,1.8,C); card(s,"FRONTEND","React compilado · rutas SPA\nURL real de API · HTTPS",8.6,3.35,3.7,1.8,O); txt(s,"Objetivo: publicar sin Docker y sin exponer secretos.",.72,6.4,11,.4,15,G,True)

s=base("PANORAMA",2,"Arquitectura recomendada","En local se usa XAMPP; en Hostinger compartido PHP y MySQL/MariaDB se ejecutan sin Docker.")
card(s,"LOCAL · XAMPP","React :5173\nApache :80\nSymfony /cine-app/public\nMariaDB :3306",.7,2.2,3.3,3.35,C); txt(s,"→",4.25,3.45,.6,.5,30,G,True,align=PP_ALIGN.CENTER)
card(s,"PRODUCCIÓN","tudominio.com\nFrontend: archivos dist\n\napi.tudominio.com\nSymfony + MySQL",5.05,2.2,3.65,3.35,G,"IDEAL"); card(s,"VENTAJAS","Rutas más claras\nDespliegues separados\nLogs independientes\nEscala mejor",9.0,2.2,3.3,3.35,O)
txt(s,"Alternativa: React dentro de Symfony public/, con fallback SPA correctamente configurado.",.72,6.28,11.4,.35,13,M)

s=base("BACKEND",3,"Estado actual y cambios necesarios","Hallazgos reales del repositorio.")
card(s,"YA ESTÁ BIEN","✓ Symfony 7.4\n✓ PHP >= 8.2\n✓ Doctrine y migraciones\n✓ Caché ORM para prod",.7,2.2,3.7,3.5,G)
card(s,"CAMBIAR","APP_ENV está en dev\nAPP_SECRET es de ejemplo\nDB usa 127.0.0.1\nMySQL está fijado en 8.4",4.8,2.2,3.7,3.5,R)
card(s,"PREPARAR",".htaccess / Apache rewrite\nCORS si hay subdominio\nPermisos de var/\nInstalación de producción",8.9,2.2,3.7,3.5,O)
txt(s,"No subir .env.local, contraseñas, dumps SQL ni logs.",.73,6.27,11.4,.35,14,R,True)

s=base("BACKEND",4,"Variables de producción","Guardar secretos en .env.local o variables del servidor.")
box(s,.72,2.05,7.2,4.05,N2,True); txt(s,'APP_ENV=prod\nAPP_DEBUG=0\nAPP_SECRET=<cadena-aleatoria-larga>\nDATABASE_URL="mysql://USER:CLAVE@HOST:3306/BASE?charset=utf8mb4"\nDEFAULT_URI=https://api.tudominio.com',1.05,2.5,6.5,3,15,W,False,"Consolas")
card(s,"BASE HOSTINGER","Copiar host, base, usuario y clave desde hPanel. Confirmar MySQL/MariaDB y ajustar serverVersion; no asumir 8.4.",8.25,2.05,4.25,1.85,O)
card(s,"SECRETO REAL","Generar APP_SECRET largo y aleatorio. Nunca reutilizar el valor de ejemplo ni publicarlo en Git.",8.25,4.25,4.25,1.85,R)

s=base("BACKEND",5,"Estructura segura de carpetas","Hostinger usa public_html y no permite cambiar libremente el document root.")
txt(s,"/home/usuario/",.85,2.1,3,.3,17,C,True,"Consolas"); txt(s,"cine-backend/       ← privado\n├── config/\n├── migrations/\n├── src/\n├── var/\n└── vendor/\n\npublic_html/api/   ← público\n├── index.php\n├── .htaccess\n└── assets",1.05,2.55,5,3.7,16,W,False,"Consolas")
card(s,"REGLA DE ORO","Sólo Symfony public/ debe quedar accesible por web. src/, config/, vendor/ y .env.local no deben exponerse.",7,2.15,5.25,1.8,R,"SEGURIDAD")
card(s,"LIMITACIÓN","Si el subdominio no apunta a public/, separar el núcleo fuera de public_html y adaptar index.php, o usar reglas Apache verificadas.",7,4.25,5.25,2,O)

s=base("BACKEND",6,"Pasos de despliegue","Hacer backup antes de cada actualización.")
steps=[("1 · hPanel","PHP 8.2+; pdo_mysql, ctype e iconv"),("2 · MySQL","Crear base y guardar credenciales"),("3 · Código","Subir por Git/SSH o archivo"),("4 · Composer","composer2 install --no-dev --optimize-autoloader"),("5 · Doctrine","Ejecutar migraciones en prod"),("6 · Verificar","Caché, permisos y /api/health")]
for i,(t,b) in enumerate(steps): card(s,t,b,.7+(i%3)*4.15,2.1+(i//3)*2.05,3.75,1.65,C if i<3 else G)
txt(s,"Sin Composer/SSH: generar vendor/ localmente y subirlo junto al backend.",.75,6.45,11.5,.35,13,M)

s=base("BACKEND",7,"Apache, CORS y seguridad","Las rutas /api/... necesitan una reescritura correcta.")
card(s,".HTACCESS","Configurar Apache para enviar rutas Symfony a index.php y probar URLs profundas.",.7,2.2,3.7,2,C)
card(s,"CORS · ARCHIVO .env / .env.local","Desarrollo:\nCORS_ALLOW_ORIGIN=http://localhost:5173\n\nProducción:\nCORS_ALLOW_ORIGIN=https://tudominio.com",4.8,2.2,3.7,2,O)
card(s,"HTTPS","Activar SSL, forzar HTTPS y actualizar DEFAULT_URI. Evitar contenido mixto.",8.9,2.2,3.7,2,G)
card(s,"NO PUBLICAR",".env.local · .git · dumps · logs · tests · claves",.7,4.65,5.75,1.35,R); card(s,"PERMISOS","var/cache y var/log escribibles por PHP; el resto, mínimo necesario.",6.85,4.65,5.75,1.35,O)

s=base("FRONTEND",8,"Qué cambia en React/Vite","En producción Vite genera archivos estáticos; no se ejecuta npm run dev.")
card(s,"DESARROLLO","npm run dev\nlocalhost:5173\nNo se sube",.7,2.2,3.4,3.3,C); txt(s,"→",4.3,3.4,.55,.5,30,O,True,align=PP_ALIGN.CENTER)
card(s,"COMPILAR","VITE_API_URL=https://api.tudominio.com/api\nnpm ci\nnpm run build",5,2.2,3.55,3.3,O); txt(s,"→",8.75,3.4,.55,.5,30,G,True,align=PP_ALIGN.CENTER)
card(s,"PUBLICAR","Subir sólo dist/*\na public_html\nNo node_modules",9.4,2.2,3.15,3.3,G)
txt(s,"Desarrollo usa http://localhost/cine-app/public/api; producción debe usar la URL HTTPS definitiva.",.73,6.27,11.5,.35,14,R,True)

s=base("FRONTEND",9,"Rutas SPA, API y caché","Tres detalles que suelen romper React en producción.")
card(s,"1 · REACT ROUTER","Fallback a /index.html para recargar rutas internas. Excluir /api si comparte dominio.",.7,2.2,3.75,3.2,O)
card(s,"2 · URL API","Usar HTTPS y dominio real. Mismo dominio: preferir /api relativo y evitar CORS.",4.8,2.2,3.75,3.2,C)
card(s,"3 · ASSETS","Caché larga para assets con hash; caché corta para index.html.",8.9,2.2,3.75,3.2,G)
txt(s,"Probar: home → navegación → recarga de ruta interna → llamada real a API.",.75,6.22,11.6,.4,15,W,True)

s=base("CAMBIOS",10,"Cambios que haría en el repositorio","Prioridad antes del primer despliegue.")
items=[("P0 · Producción","APP_ENV, APP_SECRET y DATABASE_URL seguros",R),("P0 · API URL","frontend/.env.production con dominio real",R),("P0 · Apache",".htaccess para Symfony y React Router",R),("P1 · CORS","Restringido al frontend real",O),("P1 · Base","serverVersion de Hostinger",O),("P1 · Deploy","Checklist/script repetible",O),("P2 · Logs","Errores sin debug público",G),("P2 · Higiene","Excluir cachés, logs y secretos",G)]
for i,(t,b,a) in enumerate(items): card(s,t,b,.7+(i%2)*6,2.0+(i//2)*1.08,5.6,.85,a)

s=base("VALIDACIÓN",11,"Checklist antes de abrir","La publicación termina cuando funciona el flujo completo.")
checks=["PHP 8.2+ y extensiones OK","APP_ENV=prod / debug apagado","Migraciones aplicadas","HTTPS definitivo","CORS restringido","Rutas React recargan","API responde JSON","Consulta real a MySQL","Compra/reserva probada","Backup y rollback"]
for i,t in enumerate(checks):
 x=.85+(i%2)*6; y=2.05+(i//2)*.82; box(s,x,y,.38,.38,N2,True,G); txt(s,"✓",x+.02,y+.02,.34,.25,13,G,True,align=PP_ALIGN.CENTER); txt(s,t,x+.58,y-.02,4.95,.4,15,W,i in (0,5))
box(s,.85,6.35,11.6,.55,N2,True); txt(s,"GO / NO-GO: sin secretos públicos · sin debug · flujo principal completo",1.05,6.52,11.15,.2,12,W,True,align=PP_ALIGN.CENTER)

s=base("OPERACIÓN",12,"Después del lanzamiento","Producción necesita mantenimiento, no sólo una primera subida.")
card(s,"CADA DEPLOY","Backup → código → Composer → migraciones → caché → smoke test",.7,2.2,3.75,2.7,C)
card(s,"CADA SEMANA","Revisar logs, espacio, errores 5xx, SSL y backups",4.8,2.2,3.75,2.7,O)
card(s,"SI FALLA","Restaurar versión anterior; no editar producción sin registrar el cambio",8.9,2.2,3.75,2.7,R)
txt(s,"Ensayar primero en staging con subdominio y base separados; luego repetir el procedimiento en producción.",.75,5.85,11.5,.65,19,W,True)

s=base("INTEGRACIÓN",13,"Cómo se comunican frontend y backend","React no accede directamente a MariaDB: siempre consume la API HTTP de Symfony.")
card(s,"1 · REACT","useMovies / SeatsPage\nsolicitan recursos\n\nNavegador :5173",.45,2.3,2.35,2.75,O,"FRONTEND")
txt(s,"HTTP GET\nJSON →",2.9,3.15,1.05,.7,14,G,True,align=PP_ALIGN.CENTER)
card(s,"2 · API.JS","Construye la URL:\nAPI_URL + path\n\nfetch(...) ",4.0,2.3,2.35,2.75,G)
txt(s,"/api/...\n→",6.47,3.15,.85,.7,14,C,True,align=PP_ALIGN.CENTER)
card(s,"3 · SYMFONY","Apache → index.php\nController → Repository\nDoctrine ORM",7.4,2.3,2.55,2.75,C,"BACKEND")
txt(s,"SQL\n↔",10.05,3.15,.75,.7,14,G,True,align=PP_ALIGN.CENTER)
card(s,"4 · MARIADB","pelicula\nfuncion\nasiento",10.9,2.3,1.95,2.75,G,"XAMPP")
box(s,.75,5.75,11.8,.62,N2,True); txt(s,"RESPUESTA: MariaDB → Doctrine → JsonResponse → fetch() → estado de React → pantalla",1.0,5.95,11.25,.25,13,W,True,align=PP_ALIGN.CENTER)
txt(s,"Ejemplo completo: http://localhost/cine-app/public/api/funciones",.75,6.55,11.6,.32,13,C,True,align=PP_ALIGN.CENTER)

s=base("ENDPOINTS",14,"Dónde están definidos los endpoints","La anotación #[Route(...)] del controlador determina cada URL pública.")
card(s,"GET /api/health","Archivo:\nsrc/Controller/HealthController.php\n\nComprueba que la API funciona.",.55,2.05,3.8,2.25,G)
card(s,"GET /api/funciones","Archivo:\nsrc/Controller/FuncionController.php\nMétodo: list()\n\nDevuelve películas y funciones.",4.77,2.05,3.8,2.25,C)
card(s,"GET /api/funciones/{id}/asientos","Archivo:\nsrc/Controller/FuncionController.php\nMétodo: seats()\n\nDevuelve asientos de una función.",8.98,2.05,3.8,2.25,O)
card(s,"QUIÉN LOS CONSUME","hooks/useMovies.js: api('/funciones')\npages/SeatsPage.jsx: api('/funciones/{id}/asientos')",.55,4.65,5.85,1.45,O)
card(s,"QUIÉN ARMA LA URL","frontend/src/services/api.js\nBase local: http://localhost/cine-app/public/api",6.78,4.65,6,1.45,C)
txt(s,"Con Docker:  docker compose exec backend php bin/console debug:router",.7,6.45,11.9,.35,14,G,True,"Consolas",PP_ALIGN.CENTER)

s=base("HTTP",15,"Cómo pide recursos el frontend","HTTP es el protocolo de comunicación entre el navegador y el servidor.")
card(s,"GET · LEER","Pide un recurso sin modificarlo.\n\nGET /api/funciones\nGET /api/funciones/5/asientos\n\nRespuesta habitual: 200 OK",.45,2.05,2.85,3.75,C,"ACTUAL")
card(s,"POST · CREAR","Envía datos para crear un recurso.\n\nPOST /api/funciones\nBody JSON con la nueva función\n\nRespuesta habitual: 201 Created",3.65,2.05,2.85,3.75,G,"CRUD")
card(s,"PUT · ACTUALIZAR","Reemplaza o actualiza un recurso identificado.\n\nPUT /api/funciones/5\nBody JSON actualizado\n\nRespuesta habitual: 200 o 204",6.85,2.05,2.85,3.75,O,"CRUD")
card(s,"DELETE · ELIMINAR","Solicita eliminar un recurso.\n\nDELETE /api/funciones/5\n\nRespuesta habitual:\n204 No Content",10.05,2.05,2.85,3.75,R,"CRUD")
box(s,.75,6.15,11.8,.62,N2,True); txt(s,"PETICIÓN HTTP = método + URL + headers + body opcional   |   RESPUESTA = estado + headers + JSON opcional",1.0,6.35,11.25,.25,12,W,True,align=PP_ALIGN.CENTER)
txt(s,"En Cine App hoy están implementados los GET; POST, PUT y DELETE representan la ampliación del CRUD.",.75,6.78,11.7,.25,11,M,True,align=PP_ALIGN.CENTER)

s=base("BACKEND",16,"De {id} a una entidad Funcion","La búsqueda ahora está escrita paso a paso para que no dependa de una conversión automática invisible.")
card(s,"1 · URL","GET /api/funciones/5/asientos\n\nSymfony reconoce:\n{id} = 5",.4,2.15,2.3,2.85,O)
txt(s,"→",2.77,3.17,.45,.45,24,G,True,align=PP_ALIGN.CENTER)
card(s,"2 · PARÁMETRO","public function seats(int $id)\n\nEl valor 5 entra como:\n$id = 5",3.25,2.15,2.45,2.85,C)
txt(s,"→",5.77,3.17,.45,.45,24,G,True,align=PP_ALIGN.CENTER)
card(s,"3 · REPOSITORY","$this->funcionRepository\n    ->find($id)\n\nDoctrine busca funcion.id = 5",6.25,2.15,2.7,2.85,G)
txt(s,"→",9.02,3.17,.45,.45,24,G,True,align=PP_ALIGN.CENTER)
card(s,"4 · RESULTADO","Existe: objeto Funcion\ny HTTP 200\n\nNo existe: null\ny HTTP 404",9.5,2.15,3.35,2.85,R)
box(s,.65,5.55,12,.68,N2,True); txt(s,"Después: findBy(['funcion' => $funcion]) busca los asientos cuyo funcion_id corresponde a esa función.",.9,5.78,11.5,.25,13,W,True,align=PP_ALIGN.CENTER)
txt(s,"Archivo: src/Controller/FuncionController.php  ·  Método: seats(int $id)",.7,6.58,11.9,.3,13,C,True,"Consolas",PP_ALIGN.CENTER)

s=base("DOCTRINE ORM",17,"Objetos PHP en lugar de SQL repetitivo","ORM significa Object–Relational Mapping: conecta entidades PHP con tablas de una base relacional.")
card(s,"ENTIDAD PHP","Funcion.php\nPelicula.php\nAsiento.php\n\nObjetos con propiedades y relaciones",.55,2.15,3.15,3.25,C,"CÓDIGO")
txt(s,"Doctrine\ntraduce ↔",3.83,3.15,1.15,.7,14,G,True,align=PP_ALIGN.CENTER)
card(s,"MAPEO ORM","#[ORM\\Entity]\n#[ORM\\Column]\n#[ORM\\ManyToOne]\n\nDescribe cómo se guardan los objetos",5.1,2.15,3.15,3.25,G,"ORM")
txt(s,"SQL\n↔",8.38,3.15,.7,.7,14,G,True,align=PP_ALIGN.CENTER)
card(s,"TABLAS MARIADB","funcion\npelicula\nasiento\n\nFilas, columnas y claves foráneas",9.2,2.15,3.55,3.25,O,"XAMPP")
box(s,.65,5.8,12,.62,N2,True); txt(s,"Doctrine permite consultar, insertar, actualizar y eliminar usando objetos; genera SQL y transforma cada fila en una entidad.",.9,6.0,11.5,.25,13,W,True,align=PP_ALIGN.CENTER)
txt(s,"Ventaja didáctica: el código trabaja con $funcion->getPelicula(), mientras Doctrine administra pelicula_id.",.7,6.58,11.9,.3,12,C,True,align=PP_ALIGN.CENTER)

s=base("QUERYBUILDER",18,"La consulta Doctrine y su SQL equivalente","createQueryBuilder construye la consulta paso a paso; getResult() es quien finalmente la ejecuta.")
card(s,"FUNCIONREPOSITORY.PHP","$qb = $this->createQueryBuilder('f')\n  ->join('f.pelicula', 'p')\n  ->addSelect('p')\n  ->orderBy('p.id', 'ASC')\n  ->addOrderBy('f.fechaHora', 'ASC');\n\nreturn $qb->getQuery()->getResult();",.45,2.0,5.75,3.55,C,"ORM")
card(s,"SQL PURO EQUIVALENTE","SELECT f.*, p.*\nFROM funcion AS f\nINNER JOIN pelicula AS p\n  ON f.pelicula_id = p.id\nORDER BY p.id ASC,\n         f.fecha_hora ASC;",6.55,2.0,6.3,3.55,O,"SQL")
txt(s,"f = alias de Funcion/funcion   ·   p = alias de Pelicula/pelicula   ·   JOIN une ambas tablas mediante pelicula_id",.65,5.82,12,.35,13,G,True,align=PP_ALIGN.CENTER)
box(s,.65,6.3,12,.62,N2,True); txt(s,"Con ?fecha=2026-07-30:  WHERE f.fecha_hora >= :inicio AND f.fecha_hora < :fin",.9,6.5,11.5,.25,13,W,True,"Consolas",PP_ALIGN.CENTER)

s=base("ARQUITECTURA",19,"Una responsabilidad por capa","Separar HTTP, reglas, consultas y datos hace que el código sea más legible, reutilizable y comprobable.")
card(s,"CONTROLLER","Recibe Request\nValida formato HTTP\nElige 200/400/404\nConstruye JsonResponse\n\nNo escribe consultas SQL.",.35,2.0,2.75,3.75,O,"HTTP")
card(s,"SERVICE","Coordina casos de uso\ny reglas del negocio.\n\nEjemplo actual:\nHealthService\n\nNo es obligatorio si sólo delega.",3.3,2.0,2.75,3.75,G,"NEGOCIO")
card(s,"REPOSITORY","Centraliza consultas\nQueryBuilder / Doctrine\nFiltros y ordenamiento\n\nFuncionRepository\nAsientoRepository",6.25,2.0,2.75,3.75,C,"DATOS")
card(s,"ENTITY","Representa el dominio\ny su mapeo ORM.\n\nFuncion\nPelicula\nAsiento\n\nNo conoce HTTP.",9.2,2.0,2.75,3.75,R,"MODELO")
box(s,.65,6.15,12,.62,N2,True); txt(s,"Frontend → Controller → Service (si hay reglas) → Repository → Doctrine ORM → MariaDB",.9,6.35,11.5,.25,14,W,True,align=PP_ALIGN.CENTER)

s=base("REFACTOR",20,"Cómo quedó Cine App","El Controller pide datos con métodos que expresan intención; cada Repository encapsula los detalles de acceso.")
card(s,"FUNCIONCONTROLLER.PHP","list(\n  #[MapQueryParameter]\n  ?string $fecha = null\n)\n\n$funciones = $this->funcionRepository\n  ->findCartelera($inicio, $fin);",.45,2.0,3.8,3.9,O,"CONTROLLER")
txt(s,"→",4.35,3.35,.55,.5,26,G,True,align=PP_ALIGN.CENTER)
card(s,"FUNCIONREPOSITORY.PHP","findCartelera($inicio, $fin)\nfind($id)\n\nJOIN con película\nFiltro por fecha\nOrdenamiento",5.0,2.0,3.35,3.9,C,"CONSULTAS")
txt(s,"+",8.47,3.35,.45,.5,26,G,True,align=PP_ALIGN.CENTER)
card(s,"ASIENTOREPOSITORY.PHP","findByFuncionOrdenados(\n    $funcion\n)\n\nFiltra por función\nOrdena fila y número",9.0,2.0,3.85,3.9,G,"CONSULTAS")
txt(s,"Archivos: src/Controller/FuncionController.php · src/Repository/FuncionRepository.php · src/Repository/AsientoRepository.php",.55,6.45,12.2,.35,11,C,True,"Consolas",PP_ALIGN.CENTER)

s=base("PARÁMETROS",21,"Query parameter explícito en list()","El método recibe únicamente el dato que necesita, en lugar de recibir todo Request para extraerlo manualmente.")
card(s,"URL SIN FILTRO","GET /api/funciones\n\nNo existe ?fecha\n\n$fecha = null\n\nDevuelve toda la cartelera",.5,2.15,3.45,3.45,G)
txt(s,"→",4.05,3.4,.5,.45,25,C,True,align=PP_ALIGN.CENTER)
card(s,"FIRMA EXPLÍCITA","public function list(\n  #[MapQueryParameter]\n  ?string $fecha = null\n): JsonResponse\n\nSymfony hace el mapeo",4.65,2.15,4.0,3.45,C,"SYMFONY")
txt(s,"←",8.75,3.4,.5,.45,25,C,True,align=PP_ALIGN.CENTER)
card(s,"URL CON FILTRO","GET /api/funciones\n  ?fecha=2026-07-30\n\n$fecha = '2026-07-30'\n\nValida y consulta ese día",9.35,2.15,3.45,3.45,O)
box(s,.75,6.05,11.8,.62,N2,True); txt(s,"?string significa: acepta un string o null   ·   MapQueryParameter indica que el valor viene después del ? en la URL",1.0,6.25,11.25,.25,13,W,True,align=PP_ALIGN.CENTER)

s=base("HTTP STATUS",22,"Qué significan los códigos de respuesta","El status code resume el resultado de la petición antes de que el frontend lea el JSON.")
card(s,"2xx · ÉXITO","200 OK\nConsulta exitosa con respuesta.\n\n201 Created\nRecurso creado.\n\n204 No Content\nÉxito sin body.",.4,2.05,3.0,3.95,G,"OK")
card(s,"4xx · ERROR DEL CLIENTE","400 Bad Request: datos inválidos\n401 Unauthorized: falta autenticación\n403 Forbidden: sin permiso\n404 Not Found: no existe\n405 Method Not Allowed: método incorrecto",3.65,2.05,4.15,3.95,O,"CLIENTE")
card(s,"5xx · ERROR DEL SERVIDOR","500 Internal Server Error\nFallo inesperado del backend.\n\n503 Service Unavailable\nServicio temporalmente no disponible.\n\nNo se debe mostrar debug al usuario.",8.05,2.05,4.85,3.95,R,"SERVIDOR")
box(s,.7,6.35,11.9,.55,N2,True); txt(s,"Regla: 2xx = salió bien · 4xx = corregir petición/autorización · 5xx = revisar backend, logs o infraestructura",.95,6.52,11.4,.22,12,W,True,align=PP_ALIGN.CENTER)

s=base("RESPUESTAS",23,"Ejemplos reales en Cine App","El frontend usa response.ok para distinguir 2xx de respuestas 4xx/5xx.")
card(s,"200 · CARTELERA","GET /api/funciones\n\n{\n  \"data\": [...],\n  \"total\": 4\n}\n\nLa consulta fue correcta.",.45,2.0,3.75,3.95,G,"SUCCESS")
card(s,"400 · FECHA INVÁLIDA","GET /api/funciones\n  ?fecha=2026-99-99\n\n{\n  \"error\": \"La fecha debe tener\n  formato AAAA-MM-DD.\"\n}",4.45,2.0,4.15,3.95,O,"VALIDACIÓN")
card(s,"404 · NO ENCONTRADO","GET /api/funciones/999/asientos\n\n{\n  \"error\": \"No existe una función\n  con el ID 999.\"\n}",8.85,2.0,4.05,3.95,R,"NOT FOUND")
box(s,.7,6.3,11.9,.62,N2,True); txt(s,"No devolver 200 con {error: ...}: el código HTTP debe representar correctamente el resultado para navegadores, apps y monitoreo.",.95,6.49,11.4,.27,12,W,True,align=PP_ALIGN.CENTER)

s=base("REACT",24,"Frontend separado por responsabilidades","App.jsx pasó de concentrar toda la aplicación a ocuparse únicamente del layout y las rutas.")
card(s,"APP.JSX","Compone Header y Footer\nDefine <Routes>\nConecta URL con Page\n\nNo consulta la API\nNo implementa pantallas",.4,2.0,2.85,3.95,O,"RUTAS")
card(s,"PAGES/","Pantallas completas:\nHomePage\nBillboardPage\nMovieDetailPage\nSeatsPage\nPaymentPage\nTicketsPage",3.55,2.0,2.85,3.95,C,"VISTAS")
card(s,"COMPONENTS/","Piezas reutilizables:\nHeader · Footer\nPoster · MovieTile\nFunctionCard · Stepper\nMiniMovie · Ticket\nLoading",6.7,2.0,2.85,3.95,G,"UI")
card(s,"APOYO","hooks/ → estado y efectos\nservices/ → HTTP\nutils/ → funciones puras\ndata/ → demostración\n\nCada carpeta tiene una razón.",9.85,2.0,3.05,3.95,R,"LÓGICA")
box(s,.7,6.3,11.9,.62,N2,True); txt(s,"Separar no significa crear archivos al azar: cada componente debe tener una responsabilidad clara y un nombre que exprese su intención.",.95,6.49,11.4,.27,12,W,True,align=PP_ALIGN.CENTER)

s=base("ESTRUCTURA",25,"Cómo está conformado frontend/src","Las dependencias avanzan desde las páginas hacia componentes, hooks, servicios y utilidades.")
txt(s,"src/\n├── components/\n│   ├── feedback/     Loading\n│   ├── layout/       Header, Footer\n│   ├── movies/       Poster, MovieTile, FunctionCard\n│   ├── navigation/   Stepper\n│   └── purchase/     MiniMovie, Ticket\n├── pages/            Home, Cartelera, Detalle, Asientos, Pago\n├── hooks/            useMovies\n├── services/         api.js\n├── utils/            formatters, movies\n├── data/             datos de respaldo\n└── App.jsx           layout + rutas",.7,1.85,6.4,4.95,15,W,False,"Consolas")
card(s,"COMPONENTE","Función React que recibe props y devuelve interfaz JSX. Puede reutilizarse en distintas páginas.",7.55,2.05,4.9,1.3,C)
card(s,"PAGE","Componente de nivel pantalla asociado a una ruta. Coordina componentes y estado de esa vista.",7.55,3.65,4.9,1.3,O)
card(s,"REGLA PRÁCTICA","Si una pieza aparece varias veces, tiene lógica propia o hace difícil leer la página, conviene extraerla.",7.55,5.25,4.9,1.3,G)

s=base("HOOKS",26,"Qué es un hook y qué hace useMovies()","Un hook es una función que reutiliza estado y efectos de React; por convención su nombre comienza con use.")
card(s,"1 · PÁGINA","HomePage o\nBillboardPage llama:\n\nconst { movies, loading, error }\n  = useMovies();",.45,2.15,2.75,3.2,O,"UI")
txt(s,"→",3.3,3.3,.45,.45,24,G,True,align=PP_ALIGN.CENTER)
card(s,"2 · HOOK","useMovies.js administra:\nuseState\nuseEffect\nloading y error\ndatos de respaldo",3.85,2.15,2.75,3.2,G,"ESTADO")
txt(s,"→",6.7,3.3,.45,.45,24,G,True,align=PP_ALIGN.CENTER)
card(s,"3 · SERVICE","api.js ejecuta:\n\nfetch(API_URL + path)\n\nGET /api/funciones",7.25,2.15,2.55,3.2,C,"HTTP")
txt(s,"→",9.9,3.3,.45,.45,24,G,True,align=PP_ALIGN.CENTER)
card(s,"4 · SYMFONY","Controller\nRepository\nDoctrine\nMariaDB\n\nDevuelve JSON",10.45,2.15,2.45,3.2,R,"BACKEND")
box(s,.7,5.8,11.9,.62,N2,True); txt(s,"Cuando llega el JSON, el hook actualiza movies; React vuelve a renderizar automáticamente los componentes que utilizan ese estado.",.95,5.99,11.4,.27,12,W,True,align=PP_ALIGN.CENTER)
txt(s,"El hook reutiliza lógica, no HTML: la presentación sigue viviendo en Pages y Components.",.7,6.58,11.9,.3,13,C,True,align=PP_ALIGN.CENTER)

s=base("PROBAR API",27,"Herramientas para ver y ejecutar endpoints","Las tres permiten trabajar con APIs HTTP, pero tienen objetivos y niveles de integración diferentes.")
card(s,"SWAGGER / OPENAPI","Documentación web generada desde una especificación.\n\nLista contratos, parámetros, respuestas y permite ejecutar endpoints.\n\nIdeal para documentar una API.",.4,2.0,3.75,3.5,C,"DOCUMENTA")
card(s,"POSTMAN","Aplicación independiente para crear colecciones, ambientes, pruebas y compartir peticiones.\n\nIdeal para equipos y suites grandes.",4.55,2.0,3.75,3.5,O,"COLABORA")
card(s,"THUNDER CLIENT","Extensión liviana dentro de Visual Studio Code.\n\nPermite elegir método, URL, query, headers y body, y ver status + JSON.\n\nEs la elegida para esta clase.",8.7,2.0,3.75,3.5,G,"USAREMOS")
box(s,.55,5.85,12.25,.85,N2,True); txt(s,"EJEMPLO  ·  GET  http://localhost/cine-app/public/api/funciones/1/asientos  →  Send  →  200 OK  →  { data: [...], screening: {...} }",.8,6.1,11.75,.32,12,W,True,"Consolas",PP_ALIGN.CENTER)

s=base("PERMISOS",28,"Permisos recomendados en Hostinger","Aplicar el mínimo necesario: lectura para el código, escritura sólo donde Symfony la necesita y acceso restringido a secretos.")
card(s,"755 · DIRECTORIOS","public_html/\npublic/ · assets/\nsrc/ · config/\nvendor/ · migrations/\n\nEl servidor puede recorrer y leer.",.4,2.0,2.85,3.8,C,"CARPETAS")
card(s,"644 · ARCHIVOS","*.php · *.js · *.css\nimágenes y fuentes\nindex.php · .htaccess\ncomposer.json\ncomposer.lock\n\nLectura sin escritura pública.",3.55,2.0,2.85,3.8,G,"ARCHIVOS")
card(s,"775 · ESCRITURA","var/cache/\nvar/log/\n\nSymfony necesita crear caché y logs.\n\nSi 755 funciona con el usuario PHP, preferir 755.",6.7,2.0,2.85,3.8,O,"SYMFONY")
card(s,"600 · SECRETOS",".env.local\nclaves privadas\ncertificados privados\n\nSólo el propietario puede leer y escribir.",9.85,2.0,3.05,3.8,R,"PRIVADO")
box(s,.7,6.15,11.9,.62,N2,True); txt(s,"Nunca usar chmod -R 777: permitiría que cualquier usuario o proceso modificara toda la aplicación.",.95,6.35,11.4,.25,13,W,True,align=PP_ALIGN.CENTER)

s=base("REFERENCIAS",29,"Fuentes oficiales y siguiente paso","Consultadas el 31 de julio de 2026.")
refs=[("Directorio raíz y document root","https://www.hostinger.com/support/1583494-what-is-the-path-to-your-website-s-root-home-directory-and-how-to-change-it-in-hostinger/"),("Versiones PHP en Hostinger","https://www.hostinger.com/support/4047803-how-to-change-the-php-version-for-subfolders-or-subdomains-in-hostinger/"),("Composer en hosting compartido","https://www.hostinger.com/tutorials/how-to-install-composer/"),("Crear y configurar subdominios","https://support.hostinger.com/en/articles/1583405-how-to-create-and-delete-subdomains-in-hostinger")]
for i,(label,url) in enumerate(refs):
 y=2.0+i*.78; q=box(s,.8,y,11.7,.52,N2,True); t=txt(s,label,1.05,y+.12,10.8,.24,13,W,i==0); t.text_frame.paragraphs[0].runs[0].hyperlink.address=url
txt(s,"SIGUIENTE PASO",.82,5.55,2.2,.3,12,O,True); txt(s,"Aplicar estos cambios al repositorio y generar el paquete exacto para Hostinger.",.82,5.95,11.3,.55,21,W,True)

s=base("INICIO LOCAL",30,"Cómo ejecutar el proyecto","Primero iniciar Apache y MySQL desde XAMPP Control Panel; Symfony queda servido por Apache y React por Vite.")
card(s,"1 · XAMPP / BACKEND","Iniciar:\n✓ Apache\n✓ MySQL\n\nAPI health:\nhttp://localhost/cine-app/\npublic/api/health",.4,2.0,3.75,3.85,C,"SERVIDOR")
card(s,"2 · BASE DE DATOS","Desde la raíz:\n\nC:\\xampp\\php\\php.exe\nbin\\console doctrine:database:create\n  --if-not-exists\n\nC:\\xampp\\php\\php.exe\nbin\\console doctrine:migrations:migrate",4.55,2.0,4.15,3.85,G,"PRIMERA VEZ")
card(s,"3 · FRONTEND REACT","Desde la raíz:\n\ncd frontend\nnpm install\nnpm run dev\n\nAbrir:\nhttp://localhost:5173",8.95,2.0,3.95,3.85,O,"VITE")
box(s,.7,6.2,11.9,.62,N2,True); txt(s,"Resumen: Backend = Apache + MySQL de XAMPP   |   Frontend = cd frontend && npm run dev",.95,6.4,11.4,.25,13,W,True,"Consolas",PP_ALIGN.CENTER)

s=base("ALTERNATIVA",31,"También se puede ejecutar con Docker","El código de React y Symfony no cambia: Docker Compose proporciona PHP, Apache, Node y MariaDB y sobrescribe las URLs mediante variables.")
card(s,"ARQUITECTURA","frontend\nReact + Vite :5173\n\nbackend\nSymfony + Apache :8081\n\ndatabase\nMariaDB interno",.4,2.0,3.25,3.9,C,"3 SERVICIOS")
card(s,"INICIAR","Desde la raíz:\n\ndocker compose up\n  --build -d\n\ndocker compose exec backend\n  php bin/console\n  doctrine:migrations:migrate\n  --no-interaction",3.95,2.0,4.15,3.9,G,"COMANDOS")
card(s,"ELEGIR ENTORNO","XAMPP API:\nlocalhost/cine-app/public/api\n\nDocker API:\nlocalhost:8081/api\n\nFrontend en ambos:\nlocalhost:5173\n\nNo iniciar ambos Vite a la vez.",8.4,2.0,4.5,3.9,O,"XAMPP O DOCKER")
box(s,.7,6.25,11.9,.62,N2,True); txt(s,"Sin cambios de lógica: Controllers, Repositories, entidades, componentes, páginas y hooks son los mismos.",.95,6.45,11.4,.25,13,W,True,align=PP_ALIGN.CENTER)

s=base("DOCKER",32,"Comandos útiles y acceso a MariaDB","Ejecutar desde la raíz del proyecto. Dentro de Docker no se usa C:\\xampp\\php\\php.exe.")
card(s,"SYMFONY EN BACKEND","docker compose exec backend php bin/console about\ndocker compose exec backend php bin/console debug:router\ndocker compose exec backend php bin/console doctrine:migrations:status\ndocker compose exec backend php bin/console doctrine:schema:validate\ndocker compose exec backend php bin/console cache:clear",.35,1.9,6.15,3.85,C,"CONTENEDOR")
card(s,"ENTRAR A MARIADB","docker compose exec database\n  mariadb -u cine -pcine cine\n\nLuego ejecutar:\nSHOW DATABASES;\nUSE cine;\nSHOW TABLES;\nDESCRIBE nombre_tabla;\nSELECT * FROM nombre_tabla LIMIT 20;",6.8,1.9,6.15,3.85,G,"BASE DE DATOS")
box(s,.55,6.08,12.25,.7,N2,True); txt(s,"Consulta rápida: docker compose exec database mariadb -u cine -pcine cine -e \"SHOW TABLES;\"",.8,6.3,11.75,.27,12,W,True,"Consolas",PP_ALIGN.CENTER)
txt(s,"Credenciales locales de compose.yaml: base cine · usuario cine · clave cine. No reutilizarlas en producción.",.7,6.86,11.9,.25,10,R,True,align=PP_ALIGN.CENTER)

s=base("FRONTEND",33,"De /cartelera a la pantalla","React reparte responsabilidades: cada archivo hace una parte del recorrido y una actualización de estado vuelve a dibujar la interfaz.")
card(s,"1 · APP.JSX","La URL /cartelera coincide con una Route y React monta <BillboardPage />.\n\nVIDA REAL: la cartelera de la entrada indica a qué sala debe ir el visitante.",.3,1.92,3.0,3.75,O,"MAPA")
card(s,"2 · BILLBOARD + HOOK","BillboardPage llama useMovies(fecha). El hook administra loading, movies y error.\n\nVIDA REAL: el cliente hace un pedido al empleado y espera la respuesta.",3.55,1.92,3.0,3.75,G,"PEDIDO")
card(s,"3 · API + BACKEND","api.js hace GET /api/funciones?fecha=... Symfony consulta Doctrine y MariaDB y devuelve JSON.\n\nVIDA REAL: el empleado lleva el pedido a depósito y trae la información.",6.8,1.92,3.0,3.75,C,"BÚSQUEDA")
card(s,"4 · ESTADO + COMPONENTES","setMovies guarda el resultado. React renderiza FunctionCard; elegir un horario navega a SeatsPage.\n\nVIDA REAL: se arma la cartelera y el cliente elige función y asiento.",10.05,1.92,3.0,3.75,R,"PANTALLA")
box(s,.55,6.05,12.25,.72,N2,True); txt(s,"/cartelera → App.jsx → BillboardPage → useMovies → api.js → Symfony/MariaDB → setMovies → FunctionCard → SeatsPage",.75,6.28,11.85,.28,12,W,True,"Consolas",PP_ALIGN.CENTER)
txt(s,"Si la API falla, useMovies carga frontend/src/data/movies.js y avisa que muestra datos de demostración.",.7,6.87,11.9,.25,10,O,True,align=PP_ALIGN.CENTER)

def divider(section, title, body, accent):
 s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=N
 box(s,0,0,.18,7.5,accent); txt(s,section,.75,.75,4,.3,12,accent,True)
 txt(s,"00",12.05,.72,.65,.3,11,M,True,align=PP_ALIGN.RIGHT)
 txt(s,title,.73,2.15,10.9,1.15,38,W,True); txt(s,body,.77,3.65,10.7,1.2,18,M)
 box(s,.77,5.5,2.2,.08,accent); txt(s,"CINE APP · RECORRIDO DIDÁCTICO",.77,6.65,4,.2,9,M,True)

divider("PARTE 2", "Frontend React", "Cómo se organiza la interfaz, qué son los componentes y hooks, y cómo el navegador utiliza HTTP para pedir datos.", O)
divider("PARTE 1", "Backend Symfony", "Dónde viven los endpoints, cómo se validan parámetros y cómo Repository, Doctrine ORM y MariaDB obtienen los datos.", C)
divider("PARTE 3", "Producción", "Cómo transformar el proyecto local con XAMPP en una publicación segura y verificable dentro de Hostinger compartido.", G)

# Ejemplo concreto de componente reutilizable. En el orden narrativo aparece
# después de la estructura del frontend y antes de explicar los hooks.
s=base("COMPONENTES",26,"Ejemplo simple: el componente Poster","Una pieza pequeña se define una sola vez y recibe datos distintos mediante props.")
box(s,.45,1.95,6.35,3.95,N2,True,RGBColor(42,57,80))
txt(s,"frontend/src/components/movies/Poster.jsx",.72,2.18,5.8,.3,13,C,True,"Consolas")
txt(s,"export function Poster({ movie, large = false }) {\n  return (\n    <div\n      className={`poster ${large ? 'poster-large' : ''}`}\n      style={{ '--c1': movie.colors[0],\n               '--c2': movie.colors[1] }}\n    >\n      <span>CINE MAX</span>\n      <strong>{movie.short}</strong>\n      <small>Solo en cines</small>\n    </div>\n  );\n}",.72,2.65,5.75,2.95,12,W,False,"Consolas")
card(s,"QUÉ RECIBE","movie: datos de la película.\nlarge: decide si agrega la clase poster-large.",7.15,1.95,5.7,1.25,O,"PROPS")
card(s,"DÓNDE SE REUTILIZA","Directamente: HomePage y MovieDetailPage.\n\nTambién dentro de MovieTile, FunctionCard y MiniMovie.",7.15,3.45,5.7,1.55,G,"5 USOS")
card(s,"VENTAJA","Las páginas sólo escriben <Poster movie={movie} large /> y no repiten su estructura ni sus estilos.",7.15,5.25,5.7,1.15,C,"REUTILIZAR")

# Orden narrativo: backend → frontend/HTTP → producción → referencias.
order = [
 0, 29, 30, 31,
 34, 13, 26, 20, 15, 16, 17, 18, 19, 2,
 33, 32, 23, 24, 36, 25, 7, 8, 14, 21, 22, 12,
 35, 1, 3, 4, 27, 5, 6, 9, 10, 11,
  28,
]
slide_ids = list(prs.slides._sldIdLst)
for slide_id in slide_ids:
 prs.slides._sldIdLst.remove(slide_id)
for original_index in order:
 prs.slides._sldIdLst.append(slide_ids[original_index])

# Renumera la esquina superior derecha después de mover las diapositivas.
for slide_number, slide in enumerate(prs.slides, start=1):
 if slide_number == 1:
  continue
 for shape in slide.shapes:
  if not hasattr(shape, "text_frame") or not shape.has_text_frame:
   continue
  value = shape.text.strip()
  if len(value) == 2 and value.isdigit():
   paragraph = shape.text_frame.paragraphs[0]
   if paragraph.runs:
    paragraph.runs[0].text = f"{slide_number:02d}"
   else:
    paragraph.text = f"{slide_number:02d}"
   break

prs.save(OUT); print(OUT)
